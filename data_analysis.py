import streamlit as st
import pandas as pd

import os
import docx
import pptx
from PyPDF2 import PdfReader

from langchain.text_splitter import RecursiveCharacterTextSplitter

from langchain_core.prompts import ChatPromptTemplate
from langchain.agents import AgentExecutor, create_tool_calling_agent

from langchain_community.vectorstores import FAISS
from langchain.tools.retriever import create_retriever_tool
from langchain_community.embeddings import DashScopeEmbeddings
from langchain.chat_models import init_chat_model
from langchain_experimental.tools import PythonAstREPLTool
import matplotlib.pyplot as plt
from dotenv import load_dotenv



# 设置matplotlib
plt.ioff()  # 非交互模式
# 加载环境变量
load_dotenv(override=True)
# 获取API密钥
DeepSeek_API_KEY = os.getenv("DEEPSEEK_API_KEY")
dashscope_api_key = os.getenv("DASHSCOPE_API_KEY")
# 设置环境变量
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
# 页面配置
st.set_page_config(
    page_title="智能数据分析系统",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 自定义CSS样式
st.markdown("""
<style>
    /* 主题色彩 */
    :root {
        --primary-color: #1f77b4;
        --secondary-color: #ff7f0e;
        --success-color: #2ca02c;
        --warning-color: #ff9800;
        --error-color: #d62728;
        --background-color: #f8f9fa;
    }

    /* 隐藏默认的Streamlit样式 */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}

    /* 标题样式 */
    .main-header {
        background: linear-gradient(90deg, #1f77b4, #ff7f0e);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-size: 3rem;
        font-weight: bold;
        text-align: center;
        margin-bottom: 2rem;
    }

    /* 卡片样式 */
    .info-card {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        margin: 1rem 0;
        border-left: 4px solid var(--primary-color);
    }

    .success-card {
        background: linear-gradient(135deg, #e8f5e8, #f0f8f0);
        border-left: 4px solid var(--success-color);
    }

    .warning-card {
        background: linear-gradient(135deg, #fff8e1, #fffbf0);
        border-left: 4px solid var(--warning-color);
    }

    /* 按钮样式 */
    .stButton > button {
        background: linear-gradient(45deg, #1f77b4, #2196F3);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.5rem 1rem;
        font-weight: 600;
        transition: all 0.3s ease;
        box-shadow: 0 2px 8px rgba(31, 119, 180, 0.3);
    }

    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(31, 119, 180, 0.4);
    }

    /* Tab样式 */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        background-color: #f8f9fa;
        border-radius: 10px;
        padding: 0.5rem;
    }

    .stTabs [data-baseweb="tab"] {
        height: 60px;
        background-color: white;
        border-radius: 8px;
        padding: 0 24px;
        font-weight: 600;
        border: 2px solid transparent;
        transition: all 0.3s ease;
    }

    .stTabs [aria-selected="true"] {
        background: linear-gradient(45deg, #1f77b4, #2196F3);
        color: white !important;
        border: 2px solid #1f77b4;
    }

    /* 侧边栏样式 */
    .css-1d391kg {
        background: linear-gradient(180deg, #f8f9fa, #ffffff);
    }

    /* 文件上传区域 */
    .uploadedFile {
        background: #f8f9fa;
        border: 2px dashed #1f77b4;
        border-radius: 10px;
        padding: 1rem;
        text-align: center;
        margin: 1rem 0;
    }

    /* 状态指示器 */
    .status-indicator {
        display: inline-flex;
        align-items: center;
        gap: 0.5rem;
        padding: 0.5rem 1rem;
        border-radius: 20px;
        font-weight: 600;
        font-size: 0.9rem;
    }

    .status-ready {
        background: #e8f5e8;
        color: #2ca02c;
        border: 1px solid #2ca02c;
    }

    .status-waiting {
        background: #fff8e1;
        color: #ff9800;
        border: 1px solid #ff9800;
    }

    /* 聊天界面样式 */
    .chat-container {
        display: flex;
        flex-direction: column;
        height: 70vh;
        overflow-y: auto;
        padding: 1rem;
        border-radius: 10px;
        background-color: #f8f9fa;
    }

    .chat-message {
        margin-bottom: 1rem;
        padding: 0.75rem 1rem;
        border-radius: 8px;
        max-width: 80%;
    }

    .user-message {
        background-color: #e8f5e8;
        align-self: flex-end;
    }

    .assistant-message {
        background-color: white;
        align-self: flex-start;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
    }

    .message-content {
        margin: 0;
    }

    .chat-input-container {
        margin-top: 1rem;
    }

    .image-container {
        margin-top: 0.5rem;
        border-radius: 8px;
        overflow: hidden;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
    }
</style>
""", unsafe_allow_html=True)


# 初始化embeddings
@st.cache_resource
def init_embeddings():
    return DashScopeEmbeddings(
        model="text-embedding-v1",
        dashscope_api_key=dashscope_api_key
    )


# 初始化LLM
@st.cache_resource
def init_llm():
    return init_chat_model("deepseek-chat", model_provider="deepseek")


# 初始化会话状态
def init_session_state():
    default_states = {
        'pdf_messages': [],
        'csv_messages': [],
        'df': None,
        'document_chunks': [],
        'document_metadata': [],
        'last_uploaded_files': None,
        'image_counter': 0,
        'pdf_chat_history': [],
        'csv_chat_history': []
    }
    for key, value in default_states.items():
        if key not in st.session_state:
            st.session_state[key] = value


# 文档处理函数
def process_docx(docx_file):
    doc = docx.Document(docx_file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)


def process_pptx(pptx_file):
    prs = pptx.Presentation(pptx_file)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)



def process_files(files):
    full_text = ""
    for file in files:
        file_ext = os.path.splitext(file.name)[1].lower()
        try:
            if file_ext == '.pdf':
                pdf_reader = PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        full_text += page_text
            elif file_ext == '.docx':
                full_text += process_docx(file)
            elif file_ext == '.pptx':
                full_text += process_pptx(file)
            elif file_ext == '.txt':
                full_text += file.getvalue().decode('utf-8', errors='ignore')
            else:
                st.warning(f"⚠️ 不支持的文件格式: {file.name}")
        except Exception as e:
            st.warning(f"⚠️ 处理文件 {file.name} 时出错: {str(e)}")
    return full_text


def get_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = text_splitter.split_text(text)
    return chunks


def vector_store(text_chunks):
    embeddings = init_embeddings()
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_db")


def check_database_exists():
    return os.path.exists("faiss_db") and os.path.exists("faiss_db/index.faiss")


def get_pdf_response(user_question):
    if not check_database_exists():
        return "❌ 请先上传文档并点击'Submit & Process'按钮来处理文档！"

    try:
        embeddings = init_embeddings()
        llm = init_llm()

        new_db = FAISS.load_local("faiss_db", embeddings, allow_dangerous_deserialization=True)
        retriever = new_db.as_retriever()

        prompt = ChatPromptTemplate.from_messages([
            ("system",
             """你是AI助手，请根据提供的上下文回答问题，确保提供所有细节，如果答案不在上下文中，请说"答案不在上下文中"，不要提供错误的答案"""),
            ("placeholder", "{chat_history}"),
            ("human", "{input}"),
            ("placeholder", "{agent_scratchpad}"),
        ])

        retrieval_chain = create_retriever_tool(retriever, "pdf_extractor",
                                                "This tool is to give answer to queries from the pdf")
        agent = create_tool_calling_agent(llm, [retrieval_chain], prompt)
        agent_executor = AgentExecutor(agent=agent, tools=[retrieval_chain], verbose=True)

        response = agent_executor.invoke({"input": user_question})
        return response['output']

    except Exception as e:
        return f"❌ 处理问题时出错: {str(e)}"


# CSV处理函数
def get_csv_response(query: str) -> str:
    if st.session_state.df is None:
        return "请先上传CSV文件"

    # 确保图像目录存在
    image_folder = "generated_images"
    if not os.path.exists(image_folder):
        os.makedirs(image_folder)

    # 生成唯一的图像文件名
    st.session_state.image_counter += 1
    image_path = os.path.join(image_folder, f"plot_{st.session_state.image_counter}.png")

    llm = init_llm()
    locals_dict = {'df': st.session_state.df, 'plt': plt, 'image_path': image_path}
    tools = [PythonAstREPLTool(locals=locals_dict)]

    system = f"""Given a pandas dataframe `df` answer user's query.
    Here's the output of `df.head().to_markdown()` for your reference, you have access to full dataframe as `df`:
    ```
    {st.session_state.df.head().to_markdown()}
    ```
    Give final answer as soon as you have enough data, otherwise generate code using `df` and call required tool.
    If user asks you to make a graph, save it as `{image_path}`, and output GRAPH:<graph title>.
    Example:
    ```
    plt.hist(df['Age'])
    plt.xlabel('Age')
    plt.ylabel('Count')
    plt.title('Age Histogram')
    plt.savefig('{image_path}')
    plt.close()  # 关闭图像以释放资源
    ``` output: GRAPH:Age histogram
    Query:"""

    prompt = ChatPromptTemplate.from_messages([
        ("system", system),
        ("placeholder", "{chat_history}"),
        ("human", "{input}"),
        ("placeholder", "{agent_scratchpad}"),
    ])

    agent = create_tool_calling_agent(llm=llm, tools=tools, prompt=prompt)
    agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True)

    try:
        response = agent_executor.invoke({"input": query})['output']
        # 确保所有图像都被关闭
        plt.close('all')

        # 处理图像输出
        if "GRAPH:" in str(response):
            response += f"|||{image_path}"
        return response
    except Exception as e:
        return f"❌ 分析数据时出错: {str(e)}"


def setup_chinese_fonts():
    """配置matplotlib以支持中文显示"""
    import matplotlib.font_manager as fm

    # 定义支持中文的字体列表
    chinese_fonts = ['SimHei', 'WenQuanYi Micro Hei', 'Heiti TC', 'Microsoft YaHei',
                     'SimSun', 'WenQuanYi Micro Hei', 'Heiti TC']

    # 尝试使用列表中的第一个可用字体
    font_found = False
    for font in chinese_fonts:
        try:
            plt.rcParams['font.family'] = font
            # 测试字体是否能正常显示中文
            fig, ax = plt.subplots(figsize=(1, 1))
            ax.text(0.5, 0.5, '测试', fontsize=12)
            plt.close(fig)
            font_found = True
            break
        except:
            continue

    if not font_found:
        # 如果没有找到中文字体，尝试查找系统中可用的中文字体
        try:
            # 获取系统中所有可用字体
            available_fonts = fm.findSystemFonts()
            chinese_fonts = [f for f in available_fonts if
                             any(keyword in f.lower() for keyword in ['simhei', 'heiti', 'microsoft', 'yahei'])]

            if chinese_fonts:
                # 使用找到的第一个中文字体
                font_path = chinese_fonts[0]
                font_name = fm.FontProperties(fname=font_path).get_name()
                plt.rcParams['font.family'] = font_name
                font_found = True
                print(f"使用系统中找到的中文字体: {font_name}")
        except:
            pass

    if not font_found:
        # 如果没有找到中文字体，使用系统默认字体
        print("警告: 未找到中文字体，图表中的中文可能无法正确显示")
        plt.rcParams['font.family'] = plt.rcParamsDefault['font.family']

    # 设置其他必要的参数
    plt.rcParams['axes.unicode_minus'] = False  # 解决负号显示问题
    return plt

def main():
    setup_chinese_fonts()
    init_session_state()

    # 主标题
    st.markdown('<h1 class="main-header">🤖 智能数据分析系统 </h1>', unsafe_allow_html=True)
    st.markdown(
        '<div style="text-align: center; margin-bottom: 2rem; color: #666;">集文档问答与数据分析于一体的智能助手</div>',
        unsafe_allow_html=True)

    # 创建两个主要功能的标签页
    tab1, tab2 = st.tabs(["📄 文档智能问答", "📊 CSV数据分析"])

    # 文档问答模块
    with tab1:
        col1, col2 = st.columns([2, 1])

        with col1:
            st.markdown("### 💬 与文档对话")

            # 显示数据库状态
            if check_database_exists():
                st.markdown(
                    '<div class="info-card success-card"><span class="status-indicator status-ready">✅ 文档数据库已准备就绪</span></div>',
                    unsafe_allow_html=True)
            else:
                st.markdown(
                    '<div class="info-card warning-card"><span class="status-indicator status-waiting">⚠️ 请先上传并处理文档</span></div>',
                    unsafe_allow_html=True)

            # 用户输入
            pdf_query = st.chat_input("💭 向文档提问...", disabled=not check_database_exists())
            if pdf_query:
                st.session_state.pdf_messages.append({"role": "user", "content": pdf_query})
                st.session_state.pdf_chat_history.append({"role": "user", "content": pdf_query})

                with st.spinner("🤔 AI正在分析文档..."):
                    response = get_pdf_response(pdf_query)

                st.session_state.pdf_messages.append({"role": "assistant", "content": response})
                st.session_state.pdf_chat_history.append({"role": "assistant", "content": response})

            # 聊天历史
            st.markdown("### 💬 对话历史")
            chat_container = st.container()
            with chat_container:
                for message in reversed(st.session_state.pdf_messages):
                    if message["role"] == "user":
                        st.markdown(
                            f'<div class="chat-message user-message"><p class="message-content">{message["content"]}</p></div>',
                            unsafe_allow_html=True)
                    else:
                        st.markdown(
                            f'<div class="chat-message assistant-message"><p class="message-content">{message["content"]}</p></div>',
                            unsafe_allow_html=True)

        with col2:
            st.markdown("### 📁 文档管理")

            # 文件上传
            docs = st.file_uploader(
                "📎 上传文档",
                accept_multiple_files=True,
                type=['pdf', 'docx', 'pptx', 'txt'],
                help="支持上传PDF、Word、PowerPoint和文本文件"
            )

            if docs:
                st.success(f"📄 已选择 {len(docs)} 个文件")
                for i, doc in enumerate(docs, 1):
                    st.write(f"• {doc.name}")

            # 处理按钮
            if st.button("🚀 上传并处理文档", disabled=not docs, use_container_width=True):
                with st.spinner("📊 正在处理文档..."):
                    try:
                        raw_text = process_files(docs)
                        if not raw_text.strip():
                            st.error("❌ 无法从文档中提取文本")
                            return

                        text_chunks = get_chunks(raw_text)
                        st.info(f"📝 文本已分割为 {len(text_chunks)} 个片段")

                        vector_store(text_chunks)
                        st.success("✅ 文档处理完成！")
                        st.balloons()
                        st.rerun()

                    except Exception as e:
                        st.error(f"❌ 处理文档时出错: {str(e)}")

            # 清除数据库
            if st.button("🗑️ 清除文档数据库", use_container_width=True):
                try:
                    import shutil
                    if os.path.exists("faiss_db"):
                        shutil.rmtree("faiss_db")
                    st.session_state.pdf_messages = []
                    st.session_state.pdf_chat_history = []
                    st.success("数据库已清除")
                    st.rerun()
                except Exception as e:
                    st.error(f"清除失败: {e}")

    # CSV数据分析模块
    with tab2:
        col1, col2 = st.columns([2, 1])

        with col1:
            st.markdown("### 📈 数据分析对话")

            # 显示数据状态
            if st.session_state.df is not None:
                st.markdown(
                    '<div class="info-card success-card"><span class="status-indicator status-ready">✅ 数据已加载完成</span></div>',
                    unsafe_allow_html=True)
            else:
                st.markdown(
                    '<div class="info-card warning-card"><span class="status-indicator status-waiting">⚠️ 请先上传CSV文件</span></div>',
                    unsafe_allow_html=True)

            # 用户输入
            csv_query = st.chat_input("📊 分析数据...", disabled=st.session_state.df is None)
            if csv_query:
                st.session_state.csv_messages.append({"role": "user", "content": csv_query, "type": "text"})
                st.session_state.csv_chat_history.append({"role": "user", "content": csv_query})

                with st.spinner("🔄 正在分析数据..."):
                    response = get_csv_response(csv_query)

                if isinstance(response, pd.DataFrame):
                    st.session_state.csv_messages.append(
                        {"role": "assistant", "content": response, "type": "dataframe"})
                    st.session_state.csv_chat_history.append({"role": "assistant", "content": response.to_html()})
                elif "GRAPH" in str(response):
                    # 分离图像路径和文本
                    parts = response.split("|||")
                    text = parts[0][parts[0].find("GRAPH") + 6:]
                    image_path = parts[1] if len(parts) > 1 else 'plot.png'

                    st.session_state.csv_messages.append({
                        "role": "assistant",
                        "content": text,
                        "type": "image",
                        "image_path": image_path
                    })
                    st.session_state.csv_chat_history.append(
                        {"role": "assistant", "content": f"{text}<br><img src='{image_path}' alt='分析图表'>"})
                else:
                    st.session_state.csv_messages.append({"role": "assistant", "content": response, "type": "text"})
                    st.session_state.csv_chat_history.append({"role": "assistant", "content": response})

            # 聊天历史
            st.markdown("### 💬 对话历史")
            chat_container = st.container()
            with chat_container:
                for message in reversed(st.session_state.csv_messages):
                    if message["role"] == "user":
                        st.markdown(
                            f'<div class="chat-message user-message"><p class="message-content">{message["content"]}</p></div>',
                            unsafe_allow_html=True)
                    else:
                        if message["type"] == "dataframe":
                            st.markdown(f'<div class="chat-message assistant-message">', unsafe_allow_html=True)
                            st.dataframe(message["content"])
                            st.markdown('</div>', unsafe_allow_html=True)
                        elif message["type"] == "image":
                            st.markdown(
                                f'<div class="chat-message assistant-message"><p class="message-content">{message["content"]}</p></div>',
                                unsafe_allow_html=True)
                            # 使用st.image()显示图片
                            if os.path.exists(message["image_path"]):
                                st.image(message["image_path"], caption="分析图表", use_container_width=True)
                            else:
                                st.warning(f"图片不存在: {message['image_path']}")
                        else:
                            st.markdown(
                                f'<div class="chat-message assistant-message"><p class="message-content">{message["content"]}</p></div>',
                                unsafe_allow_html=True)

        with col2:
            st.markdown("### 📊 数据管理")

            # CSV文件上传
            csv_file = st.file_uploader("📈 上传CSV文件", type='csv')
            if csv_file:
                try:
                    # 尝试多种编码读取CSV文件
                    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
                    df = None
                    for encoding in encodings:
                        try:
                            csv_file.seek(0)  # 重置文件指针
                            df = pd.read_csv(csv_file, encoding=encoding)
                            st.success(f"✅ 数据加载成功! (使用编码: {encoding})")
                            break
                        except UnicodeDecodeError:
                            continue
                    else:
                        # 尝试更健壮的chardet库
                        try:
                            import chardet
                            csv_file.seek(0)
                            raw_data = csv_file.read(10000)  # 读取前10000字节
                            result = chardet.detect(raw_data)
                            csv_file.seek(0)
                            df = pd.read_csv(csv_file, encoding=result['encoding'])
                            st.success(f"✅ 数据加载成功! (使用编码: {result['encoding']})")
                        except Exception as e:
                            raise UnicodeDecodeError("无法解析CSV文件编码")

                    if df is not None:
                        st.session_state.df = df

                        # 显示数据预览
                        with st.expander("👀 数据预览", expanded=True):
                            st.dataframe(st.session_state.df.head())
                            st.write(
                                f"📏 数据维度: {st.session_state.df.shape[0]} 行 × {st.session_state.df.shape[1]} 列")
                except Exception as e:
                    st.error(f"❌ 加载CSV文件失败: {str(e)}")

            # 数据信息
            if st.session_state.df is not None:
                if st.button("📋 显示数据信息", use_container_width=True):
                    with st.expander("📊 数据统计信息", expanded=True):
                        st.write("**基本信息:**")
                        st.text(f"行数: {st.session_state.df.shape[0]}")
                        st.text(f"列数: {st.session_state.df.shape[1]}")
                        st.write("**列名:**")
                        st.write(list(st.session_state.df.columns))
                        st.write("**数据类型:**")
                        # 修复：将dtypes转换为字符串格式显示
                        dtype_info = pd.DataFrame({
                            '列名': st.session_state.df.columns,
                            '数据类型': [str(dtype) for dtype in st.session_state.df.dtypes]
                        })
                        st.dataframe(dtype_info, use_container_width=True)

            # 清除数据
            if st.button("🗑️ 清除CSV数据", use_container_width=True):
                st.session_state.df = None
                st.session_state.csv_messages = []
                st.session_state.csv_chat_history = []
                # 清除所有图像
                image_folder = "generated_images"
                if os.path.exists(image_folder):
                    for file in os.listdir(image_folder):
                        os.remove(os.path.join(image_folder, file))
                st.success("数据已清除")
                st.rerun()

    # 底部信息
    st.markdown("---")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("**🔧 技术栈:**")
        st.markdown("• LangChain • Streamlit • FAISS • DeepSeek")
    with col2:
        st.markdown("**✨ 功能特色:**")
        st.markdown("• 多格式文档问答 • 数据可视化分析")
    with col3:
        st.markdown("**💡 使用提示:**")
        st.markdown("• 支持多文件上传 • 实时对话交互")


if __name__ == "__main__":
    main()